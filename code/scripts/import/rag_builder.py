#!/usr/bin/env python3
"""RAG向量知识库构建工具 - 从技术文档wiki构建可检索的向量数据库."""

import argparse
import hashlib
import json
import os
import sys
import time
from collections import Counter, defaultdict
from pathlib import Path

WIKI_ROOT = Path("/mnt/d/知识库wiki")
RAG_DIR = WIKI_ROOT / "rag_data"
EXTRACTED_DIR = RAG_DIR / "extracted"
CHROMA_DIR = Path("/home/eric_jia/rag_chromadb")  # Linux fs for SQLite reliability
RESULT_FILE = WIKI_ROOT / "00_目录索引" / "classification_result.json"

COLLECTION_NAME = "wiki_docs"
EMBEDDING_MODEL = "BAAI/bge-base-zh-v1.5"
CHUNK_SIZE = 800
CHUNK_OVERLAP = 80   # 减小重叠，降低跨主题 chunk 混淆
MAX_TEXT_PER_DOC = 100000  # Cap text at 100K chars per doc (electrical drawings can be 3M+)


# ── Text extraction ─────────────────────────────────────────────────────

def extract_pdf(filepath):
    import pymupdf4llm
    try:
        return pymupdf4llm.to_markdown(str(filepath))
    except Exception as e:
        print(f"  [WARN] pymupdf4llm failed for {filepath}: {e}")
        # Fallback to raw pymupdf
        import pymupdf
        try:
            doc = pymupdf.open(str(filepath))
            text = ""
            for page in doc:
                text += page.get_text() + "\n"
            doc.close()
            return text
        except Exception as e2:
            print(f"  [ERROR] pymupdf also failed: {e2}")
            return ""


def extract_docx(filepath):
    from docx import Document
    try:
        doc = Document(str(filepath))
        parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text)
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    parts.append(" | ".join(cells))
        return "\n".join(parts)
    except Exception as e:
        print(f"  [ERROR] docx extraction failed for {filepath}: {e}")
        return ""


def extract_pptx(filepath):
    from pptx import Presentation
    try:
        prs = Presentation(str(filepath))
        parts = []
        for i, slide in enumerate(prs.slides):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_texts.append(text)
            if slide_texts:
                parts.append(f"[Slide {i+1}]\n" + "\n".join(slide_texts))
        return "\n\n".join(parts)
    except Exception as e:
        print(f"  [ERROR] pptx extraction failed for {filepath}: {e}")
        return ""


def extract_xlsx(filepath):
    from openpyxl import load_workbook
    try:
        wb = load_workbook(str(filepath), read_only=True, data_only=True)
        parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    rows.append(" | ".join(cells))
            if rows:
                parts.append(f"[{sheet_name}]\n" + "\n".join(rows))
        wb.close()
        return "\n\n".join(parts)
    except Exception as e:
        print(f"  [ERROR] xlsx extraction failed for {filepath}: {e}")
        return ""


def extract_xls(filepath):
    """For .xls files, try openpyxl first (sometimes works), otherwise skip."""
    try:
        return extract_xlsx(filepath)
    except Exception:
        print(f"  [WARN] .xls not supported, skipping: {filepath}")
        return ""


EXTRACTORS = {
    '.pdf': extract_pdf,
    '.docx': extract_docx,
    '.doc': extract_docx,  # python-docx can handle some .doc files
    '.pptx': extract_pptx,
    '.ppt': extract_pptx,
    '.xlsx': extract_xlsx,
    '.xls': extract_xls,
}


def load_documents():
    """Load document list from classification_result.json."""
    with open(str(RESULT_FILE), 'r', encoding='utf-8') as f:
        data = json.load(f)
    return [d for d in data['documents'] if d.get('dest_path')]


def extract_all(force=False):
    """Extract text from all documents. Cache results to extracted/ dir."""
    EXTRACTED_DIR.mkdir(parents=True, exist_ok=True)
    documents = load_documents()

    stats = {"total": len(documents), "extracted": 0, "cached": 0, "empty": 0, "failed": 0}
    results = []

    for i, doc in enumerate(documents):
        name = doc['name']
        dest = doc['dest_path']
        ext = os.path.splitext(name)[1].lower()

        # Cache key based on file path hash
        cache_key = hashlib.md5(dest.encode()).hexdigest()
        cache_file = EXTRACTED_DIR / f"{cache_key}.json"

        if not force and cache_file.exists():
            with open(str(cache_file), 'r', encoding='utf-8') as f:
                cached = json.load(f)
            results.append(cached)
            stats["cached"] += 1
            continue

        print(f"  [{i+1}/{len(documents)}] {name}", end="", flush=True)

        extractor = EXTRACTORS.get(ext)
        if not extractor:
            print(f" ... unsupported format")
            stats["failed"] += 1
            continue

        text = extractor(dest)
        text = text.strip() if text else ""

        if not text:
            print(f" ... empty")
            stats["empty"] += 1
        else:
            print(f" ... {len(text)} chars")
            stats["extracted"] += 1

        # Parse category info
        cat_parts = doc['category'].split('/')
        top_cat = cat_parts[0]
        sub_cat = cat_parts[1] if len(cat_parts) > 1 else ""

        entry = {
            "source": dest,
            "filename": name,
            "category": top_cat,
            "subcategory": sub_cat,
            "file_type": ext,
            "size": doc['size'],
            "text": text,
        }

        # Save cache
        with open(str(cache_file), 'w', encoding='utf-8') as f:
            json.dump(entry, f, ensure_ascii=False)

        results.append(entry)

    print(f"\n提取完成: 总计 {stats['total']}, 新提取 {stats['extracted']}, "
          f"缓存 {stats['cached']}, 空文件 {stats['empty']}, 失败 {stats['failed']}")
    return results


# ── Chunking ─────────────────────────────────────────────────────────────

def _preprocess_alarm_boundaries(text: str) -> str:
    """在 FANUC 报警代码前插入强力分段标记，确保现象+对策不被切散."""
    import re
    # 匹配报警代码行：SRVO-023 / SRVO － 023 / SVGN-382 等
    alarm_pat = re.compile(
        r'(?<=\n)([A-Z]{2,6}\s*[－\-]\s*\d{3,4}[A-Z]?(?:\s+\w+)?(?:\s*\([^)]*\))?\s*\n)',
        re.MULTILINE
    )
    # 每个报警代码前插入强分隔
    text = alarm_pat.sub(r'\n\n===ALARM_BOUNDARY===\n\n\1', text)
    return text


def chunk_documents(extracted_docs):
    """Split extracted text into chunks for embedding."""
    from langchain_text_splitters import RecursiveCharacterTextSplitter

    splitter = RecursiveCharacterTextSplitter(
        chunk_size=CHUNK_SIZE,
        chunk_overlap=CHUNK_OVERLAP,
        separators=[
            "===ALARM_BOUNDARY===",  # 报警代码强边界（预处理注入）
            "\n\n",                   # 段落边界
            "。", ". ", "！", "？",   # 句子边界（中英文句号/感叹/问号优先）
            "；", "\n",               # 分句/换行
            " ",                      # 词边界
            "",                       # 字符边界（兜底）
        ],
    )

    all_chunks = []
    for doc in extracted_docs:
        text = doc.get("text", "")
        if not text or len(text) < 10:
            continue

        # Cap text length to avoid excessive chunks from large electrical drawings
        if len(text) > MAX_TEXT_PER_DOC:
            text = text[:MAX_TEXT_PER_DOC]

        # 预处理：报警代码边界注入分隔标记
        text = _preprocess_alarm_boundaries(text)

        chunks = splitter.split_text(text)
        for j, chunk in enumerate(chunks):
            all_chunks.append({
                "id": f"{hashlib.md5(doc['source'].encode()).hexdigest()}_{j}",
                "text": chunk,
                "metadata": {
                    "source": doc["source"],
                    "filename": doc["filename"],
                    "category": doc["category"],
                    "subcategory": doc["subcategory"],
                    "file_type": doc["file_type"],
                    "chunk_index": j,
                    "total_chunks": len(chunks),
                },
            })

    print(f"分块完成: {len(extracted_docs)} 文档 → {len(all_chunks)} 个chunks")
    return all_chunks


# ── Vector DB ────────────────────────────────────────────────────────────

def get_embedding_function():
    """Create ChromaDB-compatible embedding function using sentence-transformers."""
    import torch
    from chromadb.utils.embedding_functions import SentenceTransformerEmbeddingFunction
    device = "cuda" if torch.cuda.is_available() else "cpu"
    print(f"  嵌入模型设备: {device}")
    return SentenceTransformerEmbeddingFunction(
        model_name=EMBEDDING_MODEL,
        device=device,
        trust_remote_code=False,
    )


def generate_auto_summaries(chunks):
    """P3: 为每个文档生成AI摘要chunk，提升检索覆盖率。

    用LLM对每个文档的前几个chunk生成100-200字摘要，
    作为特殊类型的chunk存入向量库。
    """
    try:
        from openai import OpenAI
    except ImportError:
        print("  跳过自动摘要 (openai未安装)")
        return chunks

    # 按文件分组
    docs = {}
    for c in chunks:
        fn = c["metadata"]["filename"]
        if fn not in docs:
            docs[fn] = []
        docs[fn].append(c)

    # 尝试连接LLM
    try:
        import sys
        sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
        from rag_core import MIOFFICE_API_BASE, MIOFFICE_API_KEY
        client = OpenAI(base_url=MIOFFICE_API_BASE, api_key=MIOFFICE_API_KEY, timeout=30)
    except Exception:
        print("  跳过自动摘要 (LLM连接失败)")
        return chunks

    summary_chunks = []
    for fn, doc_chunks in docs.items():
        # 取前3个chunk作为摘要素材
        sample_text = "\n".join(c["text"][:500] for c in doc_chunks[:3])
        if len(sample_text) < 50:
            continue

        try:
            resp = client.chat.completions.create(
                model="xiaomi/mimo-v2-flash",
                messages=[
                    {"role": "system", "content": "你是工业自动化文档摘要助手。用100-200字概括文档核心内容，包含关键型号、报警代码、操作步骤等重要信息。"},
                    {"role": "user", "content": f"请为以下文档生成摘要:\n\n{sample_text[:2000]}"},
                ],
                max_tokens=300,
                temperature=0.3,
            )
            summary = resp.choices[0].message.content.strip()
            if len(summary) < 20:
                continue

            summary_chunks.append({
                "id": f"summary_{hashlib.md5(fn.encode()).hexdigest()}",
                "text": f"[文档摘要] {fn}\n{summary}",
                "metadata": {
                    "source": doc_chunks[0]["metadata"]["source"],
                    "filename": fn,
                    "category": doc_chunks[0]["metadata"]["category"],
                    "subcategory": doc_chunks[0]["metadata"].get("subcategory", ""),
                    "file_type": doc_chunks[0]["metadata"]["file_type"],
                    "chunk_index": -1,
                    "total_chunks": len(doc_chunks),
                    "chunk_type": "summary",
                },
            })
            print(f"  ✅ 摘要: {fn[:40]}... ({len(summary)}字)")
        except Exception as e:
            print(f"  ⚠️ 摘要失败: {fn[:30]}... {e}")
            continue

    if summary_chunks:
        chunks = summary_chunks + chunks
        print(f"自动摘要完成: {len(summary_chunks)} 个摘要chunk已添加")
    return chunks


def build_vectordb(chunks):
    """Build ChromaDB vector database from chunks."""
    import chromadb

    CHROMA_DIR.mkdir(parents=True, exist_ok=True)
    client = chromadb.PersistentClient(path=str(CHROMA_DIR))

    # Delete existing collection if any
    try:
        client.delete_collection(COLLECTION_NAME)
        print("已删除旧的collection")
    except Exception:
        pass

    ef = get_embedding_function()
    collection = client.get_or_create_collection(
        name=COLLECTION_NAME,
        embedding_function=ef,
        metadata={"hnsw:space": "cosine"},
    )

    # Batch insert
    batch_size = 256
    total = len(chunks)
    print(f"开始向量化入库: {total} chunks, batch_size={batch_size}")

    for start in range(0, total, batch_size):
        end = min(start + batch_size, total)
        batch = chunks[start:end]

        ids = [c["id"] for c in batch]
        documents = [c["text"] for c in batch]
        metadatas = [c["metadata"] for c in batch]

        collection.add(ids=ids, documents=documents, metadatas=metadatas)
        print(f"  [{end}/{total}] ({end*100//total}%)", flush=True)

    print(f"向量库构建完成: {collection.count()} vectors in '{COLLECTION_NAME}'")

    # M3: 使 BM25 磁盘缓存失效，避免下次服务启动加载旧索引
    _bm25_pkl = CHROMA_DIR / "bm25_index.pkl"
    if _bm25_pkl.exists():
        try:
            _bm25_pkl.unlink()
            print("  BM25 缓存已失效 (bm25_index.pkl)")
        except OSError as e:
            print(f"  警告: BM25 缓存清理失败: {e}")

    # M3: 使 BM25/实体索引失效, 下次检索自动重建 (评审 F1)
    try:
        from rag_core import invalidate_indexes
        invalidate_indexes()
        print("  内存索引已失效 (invalidate_indexes)")
    except Exception as e:
        print(f"  警告: invalidate_indexes 调用失败: {e}")

    return collection


# ── Query ────────────────────────────────────────────────────────────────

def query(text, top_k=5):
    """Query the vector database."""
    import chromadb

    client = chromadb.PersistentClient(path=str(CHROMA_DIR))
    ef = get_embedding_function()
    collection = client.get_collection(name=COLLECTION_NAME, embedding_function=ef)

    results = collection.query(query_texts=[text], n_results=top_k)

    print(f"\n查询: \"{text}\"")
    print(f"返回 {len(results['documents'][0])} 条结果:\n")

    for i, (doc, meta, dist) in enumerate(zip(
        results['documents'][0],
        results['metadatas'][0],
        results['distances'][0],
    )):
        score = 1 - dist  # cosine similarity
        print(f"{'='*60}")
        print(f"[{i+1}] 相关度: {score:.4f}")
        print(f"    文件: {meta['filename']}")
        print(f"    分类: {meta['category']}/{meta['subcategory']}")
        print(f"    块:   {meta['chunk_index']+1}/{meta['total_chunks']}")
        print(f"    路径: {meta['source']}")
        print(f"{'─'*60}")
        # Truncate long chunks for display
        display = doc[:300] + "..." if len(doc) > 300 else doc
        print(f"    {display}")
        print()

    return results


# ── Knowledge points ─────────────────────────────────────────────────────

def generate_knowledge_points(extracted_docs):
    """Generate knowledge point summary by category."""
    import re

    # Chinese stop words (minimal set)
    stop_words = set("的了是在不有和人这中大为上个国我以要他时来用们生到作地于出会"
                     "可也你对就说等都而及与或其它被从已将能使把那但又如此因之最所"
                     "更些什么什麼下之没有能时候如果因为所以但是虽然不过只是这个那个"
                     "一个两个三个已经正在可以比较")

    def extract_keywords(text, top_n=20):
        # Simple word frequency (Chinese + English)
        words = re.findall(r'[\u4e00-\u9fff]{2,}|[A-Za-z][A-Za-z0-9_]{2,}', text)
        words = [w for w in words if w not in stop_words and len(w) >= 2]
        counter = Counter(words)
        return [{"word": w, "count": c} for w, c in counter.most_common(top_n)]

    by_category = defaultdict(lambda: defaultdict(list))
    for doc in extracted_docs:
        if doc.get("text") and len(doc["text"]) >= 10:
            by_category[doc["category"]][doc["subcategory"]].append(doc)

    knowledge = {}
    for cat in sorted(by_category.keys()):
        knowledge[cat] = {}
        for sub in sorted(by_category[cat].keys()):
            docs = by_category[cat][sub]
            all_text = " ".join(d["text"][:2000] for d in docs)  # Sample for keywords
            keywords = extract_keywords(all_text)
            knowledge[cat][sub] = {
                "doc_count": len(docs),
                "total_chars": sum(len(d["text"]) for d in docs),
                "keywords": keywords,
                "sample_files": [d["filename"] for d in docs[:5]],
            }

    # Save JSON
    RAG_DIR.mkdir(parents=True, exist_ok=True)
    kp_json = RAG_DIR / "knowledge_points.json"
    with open(str(kp_json), 'w', encoding='utf-8') as f:
        json.dump(knowledge, f, ensure_ascii=False, indent=2)

    # Save Markdown
    kp_md = RAG_DIR / "knowledge_points.md"
    with open(str(kp_md), 'w', encoding='utf-8') as f:
        f.write("# 知识点提取归档\n\n")
        f.write(f"> 生成时间: {time.strftime('%Y-%m-%d %H:%M')}\n\n")

        total_docs = sum(
            info["doc_count"]
            for cat in knowledge.values()
            for info in cat.values()
        )
        total_chars = sum(
            info["total_chars"]
            for cat in knowledge.values()
            for info in cat.values()
        )
        f.write(f"共 **{total_docs}** 份有效文档, 提取文本 **{total_chars:,}** 字符\n\n")

        for cat in sorted(knowledge.keys()):
            f.write(f"## {cat}\n\n")
            for sub in sorted(knowledge[cat].keys()):
                info = knowledge[cat][sub]
                label = sub if sub else "(全部)"
                f.write(f"### {label}\n\n")
                f.write(f"- 文档数: {info['doc_count']}\n")
                f.write(f"- 文本量: {info['total_chars']:,} 字符\n")
                f.write(f"- 关键词: {', '.join(kw['word'] for kw in info['keywords'][:10])}\n")
                f.write(f"- 示例文件: {', '.join(info['sample_files'])}\n\n")

    print(f"知识点归档已生成: {kp_json}, {kp_md}")
    return knowledge


# ── Stats ────────────────────────────────────────────────────────────────

def show_stats():
    """Show vector database statistics."""
    import chromadb

    if not CHROMA_DIR.exists():
        print("向量库尚未构建，请先运行 build")
        return

    client = chromadb.PersistentClient(path=str(CHROMA_DIR))
    try:
        collection = client.get_collection(name=COLLECTION_NAME)
    except Exception:
        print("Collection不存在，请先运行 build")
        return

    count = collection.count()
    print(f"向量库统计:")
    print(f"  Collection: {COLLECTION_NAME}")
    print(f"  向量数量: {count}")
    print(f"  嵌入模型: {EMBEDDING_MODEL}")
    print(f"  存储路径: {CHROMA_DIR}")

    # Count extracted files
    if EXTRACTED_DIR.exists():
        extracted = len(list(EXTRACTED_DIR.glob("*.json")))
        print(f"  已提取文档: {extracted}")

    # Knowledge points
    kp = RAG_DIR / "knowledge_points.json"
    if kp.exists():
        with open(str(kp), 'r', encoding='utf-8') as f:
            data = json.load(f)
        cats = len(data)
        subs = sum(len(v) for v in data.values())
        print(f"  知识点分类: {cats} 大类, {subs} 子类")


# ── Main ─────────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(description="RAG向量知识库构建工具")
    sub = parser.add_subparsers(dest="command")

    sub.add_parser("extract", help="仅提取文本")
    sub.add_parser("build", help="全流程: 提取+分块+向量化+知识点归档")
    sub.add_parser("stats", help="显示库统计信息")

    q = sub.add_parser("query", help="查询向量库")
    q.add_argument("text", help="查询文本")
    q.add_argument("--top_k", type=int, default=5, help="返回结果数 (default: 5)")

    args = parser.parse_args()

    if args.command == "extract":
        extract_all()

    elif args.command == "build":
        print("=" * 60)
        print("RAG向量知识库构建")
        print("=" * 60)

        print("\n[1/4] 提取文本...")
        docs = extract_all()

        print("\n[2/4] 文本分块...")
        chunks = chunk_documents(docs)

        # P3: 自动摘要入库
        print("\n[2.5/4] 自动生成文档摘要...")
        chunks = generate_auto_summaries(chunks)

        print("\n[3/4] 向量化入库...")
        build_vectordb(chunks)

        print("\n[4/4] 知识点归档...")
        generate_knowledge_points(docs)

        print("\n" + "=" * 60)
        print("全部完成!")
        show_stats()

    elif args.command == "query":
        query(args.text, args.top_k)

    elif args.command == "stats":
        show_stats()

    else:
        parser.print_help()


if __name__ == "__main__":
    main()
