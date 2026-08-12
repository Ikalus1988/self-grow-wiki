#!/usr/bin/env python3
"""批量导入工具 — 从指定目录导入 PDF/DOCX/PPTX/XLSX 到 ChromaDB.

用法:
  python3 import_batch.py /mnt/c/Users/hp/Downloads/待入库
  python3 import_batch.py /mnt/c/Users/hp/Downloads/待入库 --dry-run
  python3 import_batch.py /mnt/c/Users/hp/Downloads/待入库 --category "FANUC机器人"
"""

import os
import sys
import json
import hashlib
import logging
import argparse
from pathlib import Path
from typing import List, Dict

import chromadb
from chromadb.config import Settings
from chromadb.utils.embedding_functions import SentenceTransformerEmbeddingFunction
from langchain_text_splitters import RecursiveCharacterTextSplitter

logging.basicConfig(level=logging.INFO, format="%(asctime)s [%(levelname)s] %(message)s")
log = logging.getLogger("import_batch")

# ── 配置 ──────────────────────────────────────────────────────────────

CHROMA_DIR = Path("/home/eric_jia/rag_chromadb")
COLLECTION_NAME = "wiki_docs"
EMBEDDING_MODEL = "BAAI/bge-base-zh-v1.5"

CHUNK_SIZE = 800
CHUNK_OVERLAP = 100
CHUNK_SEPARATORS = ["\n\n", "\n", "。", "；", " ", ""]
MAX_TEXT_LEN = 100000
BATCH_SIZE = 256

SUPPORTED_EXTS = {".pdf", ".PDF", ".docx", ".DOCX", ".pptx", ".PPTX", ".xlsx", ".XLSX"}


# ── 文本提取 ──────────────────────────────────────────────────────────

def extract_pdf(filepath: str) -> str:
    try:
        import pymupdf4llm
        return pymupdf4llm.to_markdown(filepath)
    except Exception as e:
        log.warning("pymupdf4llm 失败: %s, 回退到 pymupdf", e)
        import fitz
        doc = fitz.open(filepath)
        text = "\n\n".join(page.get_text() for page in doc)
        doc.close()
        return text


def extract_docx(filepath: str) -> str:
    from docx import Document
    doc = Document(filepath)
    return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())


def extract_pptx(filepath: str) -> str:
    from pptx import Presentation
    prs = Presentation(filepath)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text)
    return "\n\n".join(texts)


def extract_xlsx(filepath: str) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(filepath, data_only=True)
    texts = []
    for ws in wb.worksheets:
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                texts.append(" | ".join(cells))
    return "\n".join(texts)


EXTRACTORS = {
    ".pdf": extract_pdf, ".PDF": extract_pdf,
    ".docx": extract_docx, ".DOCX": extract_docx,
    ".pptx": extract_pptx, ".PPTX": extract_pptx,
    ".xlsx": extract_xlsx, ".XLSX": extract_xlsx,
}


# ── 分块 ──────────────────────────────────────────────────────────────

def chunk_text(text: str, source: str, filename: str, category: str) -> List[Dict]:
    if len(text) > MAX_TEXT_LEN:
        text = text[:MAX_TEXT_LEN]

    splitter = RecursiveCharacterTextSplitter(
        chunk_size=CHUNK_SIZE,
        chunk_overlap=CHUNK_OVERLAP,
        separators=CHUNK_SEPARATORS,
    )
    parts = splitter.split_text(text)

    chunks = []
    src_hash = hashlib.md5(source.encode()).hexdigest()
    for i, part in enumerate(parts):
        chunks.append({
            "id": f"{src_hash}_{i}",
            "text": part,
            "metadata": {
                "source": source,
                "filename": filename,
                "category": category,
                "subcategory": "",
                "file_type": Path(filename).suffix.lstrip(".").lower(),
                "chunk_index": i,
                "total_chunks": len(parts),
            },
        })
    return chunks


# ── 主流程 ─────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(description="批量导入文档到 ChromaDB")
    parser.add_argument("source_dir", help="包含 PDF/DOCX 的目录")
    parser.add_argument("--dry-run", action="store_true", help="只显示计划不执行")
    parser.add_argument("--category", default="FANUC机器人", help="分类标签")
    args = parser.parse_args()

    source_dir = Path(args.source_dir)
    if not source_dir.exists():
        print(f"错误: 目录不存在 {source_dir}")
        sys.exit(1)

    # 扫描文件
    files = []
    for ext in SUPPORTED_EXTS:
        files.extend(source_dir.glob(f"*{ext}"))
    files = sorted(set(files))

    if not files:
        print(f"未找到支持的文件 (支持: {', '.join(SUPPORTED_EXTS)})")
        sys.exit(0)

    print(f"\n找到 {len(files)} 个文件:")
    for f in files:
        size_mb = f.stat().st_size / 1024 / 1024
        print(f"  {f.name} ({size_mb:.1f} MB)")

    if args.dry_run:
        print("\n--dry-run 模式，不执行导入")
        return

    # 初始化 ChromaDB
    print("\n初始化 ChromaDB...")
    ef = SentenceTransformerEmbeddingFunction(
        model_name=EMBEDDING_MODEL,
        device="cpu",
    )
    client = chromadb.PersistentClient(
        path=str(CHROMA_DIR),
        settings=Settings(anonymized_telemetry=False),
    )
    collection = client.get_or_create_collection(
        name=COLLECTION_NAME,
        embedding_function=ef,
        metadata={"hnsw:space": "cosine"},
    )
    print(f"  当前向量数: {collection.count()}")

    # 提取 + 分块 + 入库
    all_chunks = []
    stats = {"ok": 0, "skip": 0, "fail": 0}

    for i, f in enumerate(files):
        print(f"\n[{i+1}/{len(files)}] {f.name}", end="", flush=True)
        ext = f.suffix
        extractor = EXTRACTORS.get(ext)
        if not extractor:
            print(" ... 不支持的格式")
            stats["fail"] += 1
            continue

        try:
            text = extractor(str(f))
            text = text.strip() if text else ""
            if not text:
                print(" ... 空文件")
                stats["skip"] += 1
                continue

            chunks = chunk_text(text, str(f), f.name, args.category)
            print(f" ... {len(text)} chars → {len(chunks)} chunks")
            all_chunks.extend(chunks)
            stats["ok"] += 1

        except Exception as e:
            print(f" ... 失败: {e}")
            stats["fail"] += 1
            continue

    if not all_chunks:
        print("\n没有可导入的内容")
        return

    # 批量 upsert
    print(f"\n正在写入 {len(all_chunks)} 个 chunks...")
    for start in range(0, len(all_chunks), BATCH_SIZE):
        batch = all_chunks[start:start + BATCH_SIZE]
        collection.upsert(
            ids=[c["id"] for c in batch],
            documents=[c["text"] for c in batch],
            metadatas=[c["metadata"] for c in batch],
        )
        print(f"  {min(start + BATCH_SIZE, len(all_chunks))}/{len(all_chunks)}")

    print(f"\n完成! 导入 {stats['ok']} 个文件, 跳过 {stats['skip']}, 失败 {stats['fail']}")
    print(f"向量库总量: {collection.count()}")


if __name__ == "__main__":
    main()
